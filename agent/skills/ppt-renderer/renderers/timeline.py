from __future__ import annotations

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from .base import RenderContext, content_items


def render_timeline(ctx: RenderContext, design_slide: dict) -> None:
    content = ctx.content_for(design_slide)
    slide = ctx.blank_slide()
    ctx.add_title(slide, content["title"], design_slide["slide_number"])
    ctx.add_key_message(slide, content["key_message"])
    events = content_items(content)[:7]
    left, right, y = 1.0, 12.25, 3.55
    if len(events) > 1:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left), Inches(y), Inches(right), Inches(y))
        line.line.color.rgb = ctx.color("divider")
        line.line.width = Inches(0.03)
    gap = (right - left) / max(len(events) - 1, 1)
    for index, event in enumerate(events):
        x = left + index * gap
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.18), Inches(y - 0.18), Inches(0.36), Inches(0.36))
        node.fill.solid(); node.fill.fore_color.rgb = ctx.color("accent_primary" if index < len(events) - 1 else "accent_secondary"); node.line.fill.background()
        ctx.add_text(slide, f"{index + 1:02d}", x - 0.35, 2.65, 0.7, 0.28, size=11, color="accent_primary", bold=True, align=PP_ALIGN.CENTER)
        ctx.add_text(slide, event, max(0.58, x - 0.82), 4.05, 1.65, 1.15, size=14, color="text_primary", bold=True, align=PP_ALIGN.CENTER)
    ctx.add_footer(slide)
