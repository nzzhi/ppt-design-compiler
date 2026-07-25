from __future__ import annotations

from dataclasses import dataclass
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE


FALLBACK_PALETTE = {
    "background": "FBFCFE",
    "surface": "EEF3F8",
    "text_primary": "111827",
    "text_secondary": "4B5563",
    "accent_primary": "315C8C",
    "accent_secondary": "6A7D39",
    "warning": "B45309",
    "divider": "D6DEE8",
}


@dataclass
class RenderContext:
    presentation: Any
    design_plan: dict[str, Any]
    content_plan: dict[str, Any]
    design_system: dict[str, Any]

    def __post_init__(self) -> None:
        self.content_slides = {
            slide["slide_id"]: slide for slide in self.content_plan["slides"]
        }
        palette = self.design_system.get("palette", {})
        self.palette = {
            key: str(palette.get(key, value)).lstrip("#")
            for key, value in FALLBACK_PALETTE.items()
        }
        typography = self.design_system.get("typography", {})
        self.title_font = typography.get("title_font", "Microsoft YaHei")
        self.body_font = typography.get("body_font", "Microsoft YaHei")

    def content_for(self, design_slide: dict[str, Any]) -> dict[str, Any]:
        return self.content_slides[design_slide["slide_id"]]

    def color(self, token: str) -> RGBColor:
        value = self.palette[token]
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    def blank_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.color("background")
        return slide

    def add_text(
        self,
        slide,
        text: str,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        size: float,
        color: str = "text_primary",
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        font: str | None = None,
        margin: float = 0.02,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        run = paragraph.runs[0]
        run.font.name = font or self.body_font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = self.color(color)
        return box

    def add_title(self, slide, title: str, slide_number: int) -> None:
        size = 30 if len(title) <= 24 else 27
        self.add_text(
            slide,
            title,
            0.72,
            0.42,
            11.45,
            0.72,
            size=size,
            bold=True,
            font=self.title_font,
        )
        self.add_text(
            slide,
            f"{slide_number:02d}",
            12.15,
            0.48,
            0.48,
            0.32,
            size=10,
            color="text_secondary",
            align=PP_ALIGN.RIGHT,
        )

    def add_key_message(self, slide, message: str, y: float = 1.32) -> None:
        self.add_text(
            slide,
            message,
            0.76,
            y,
            11.75,
            0.56,
            size=16,
            color="accent_primary",
            bold=True,
        )

    def add_footer(self, slide, label: str = "PPT AGENT") -> None:
        line = slide.shapes.add_shape(1, Inches(0.72), Inches(7.06), Inches(11.9), Inches(0.012))
        line.fill.solid()
        line.fill.fore_color.rgb = self.color("divider")
        line.line.fill.background()
        self.add_text(slide, label, 0.72, 7.12, 2.0, 0.2, size=8, color="text_secondary")

    def add_panel(self, slide, x: float, y: float, width: float, height: float, *, fill: str = "surface", line: str | None = None):
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self.color(fill)
        if line:
            panel.line.color.rgb = self.color(line)
        else:
            panel.line.fill.background()
        return panel

    def add_accent_rule(self, slide, x: float, y: float, width: float, color: str = "accent_primary"):
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.06))
        rule.fill.solid()
        rule.fill.fore_color.rgb = self.color(color)
        rule.line.fill.background()
        return rule


def content_items(slide_data: dict[str, Any]) -> list[str]:
    items: list[str] = []
    for block in slide_data.get("content_blocks", []):
        content = block.get("content", "")
        if isinstance(content, list):
            items.extend(str(item) for item in content)
        elif isinstance(content, dict):
            items.extend(f"{key}: {value}" for key, value in content.items())
        elif content:
            items.append(str(content))
    return items


def split_comparison(slide_data: dict[str, Any]) -> tuple[list[str], list[str]]:
    blocks = slide_data.get("content_blocks", [])
    left = _as_items(blocks[0].get("content", [])) if blocks else []
    right = _as_items(blocks[1].get("content", [])) if len(blocks) > 1 else []
    return left, right


def _as_items(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item) for item in value]
    return [str(value)] if value else []
