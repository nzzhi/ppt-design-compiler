from __future__ import annotations

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from .base import RenderContext, content_items


def render_process(ctx: RenderContext, design_slide: dict) -> None:
    content = ctx.content_for(design_slide)
    slide = ctx.blank_slide()
    ctx.add_title(slide, content["title"], design_slide["slide_number"])
    ctx.add_key_message(slide, content["key_message"])
    steps = content_items(content)[:6]
    count = max(len(steps), 1)
    left = 0.9
    right = 12.45
    center_y = 3.55
    spacing = (right - left) / max(count - 1, 1)

    if count > 1:
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(left),
            Inches(center_y),
            Inches(right),
            Inches(center_y),
        )
        connector.line.color.rgb = ctx.color("divider")
        connector.line.width = Inches(0.025)

    for index, step in enumerate(steps):
        x = left + index * spacing
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - 0.31),
            Inches(center_y - 0.31),
            Inches(0.62),
            Inches(0.62),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = ctx.color("accent_primary" if index < count - 1 else "accent_secondary")
        node.line.fill.background()
        ctx.add_text(slide, str(index + 1), x - 0.3, center_y - 0.24, 0.6, 0.4, size=15, color="background", bold=True, align=PP_ALIGN.CENTER)
        ctx.add_text(slide, step, max(x - 1.18, 0.55), 4.12, 2.36, 1.35, size=15, color="text_primary", bold=True, align=PP_ALIGN.CENTER)

    ctx.add_text(slide, "从左到右推进", 0.92, 2.55, 2.5, 0.32, size=10, color="text_secondary")
    ctx.add_footer(slide)
