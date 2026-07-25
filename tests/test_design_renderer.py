import importlib.util
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
RENDERER_PATH = ROOT / "agent" / "skills" / "ppt-renderer" / "renderer.py"
SPEC = importlib.util.spec_from_file_location("ppt_agent_renderer", RENDERER_PATH)
MODULE = importlib.util.module_from_spec(SPEC)
assert SPEC.loader is not None
SPEC.loader.exec_module(MODULE)


class DesignRendererTests(unittest.TestCase):
    def test_design_plan_renders_all_slides_with_distinct_silhouettes(self):
        design_plan = ROOT / "projects" / "ai-drug-discovery-course-report" / "plan" / "design-plan.json"
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "design-plan.pptx"
            MODULE.render_pptx(design_plan, output)
            presentation = Presentation(output)

        self.assertEqual(len(presentation.slides), 10)
        shape_profiles = {
            tuple(shape.shape_type for shape in slide.shapes)
            for slide in presentation.slides
        }
        self.assertGreaterEqual(len(shape_profiles), 4)

    def test_renderer_dispatches_core_page_types(self):
        self.assertIn("cover", MODULE.RENDERERS)
        self.assertIn("comparison", MODULE.RENDERERS)
        self.assertIn("process", MODULE.RENDERERS)
        self.assertIn("closing", MODULE.RENDERERS)

    def test_renderer_dispatches_design_visuals_to_components(self):
        self.assertIs(MODULE._renderer_for({"slide_type": "timeline"}), MODULE.render_timeline)
        self.assertIs(MODULE._renderer_for({"slide_type": "data-insight", "visual_type": "chart"}), MODULE.render_data)
        self.assertIs(MODULE._renderer_for({"slide_type": "claim-evidence", "visual_type": "photo"}), MODULE.render_image_story)
        self.assertIs(MODULE._renderer_for({"slide_type": "claim-evidence", "visual_type": "metric"}), MODULE.render_data)


if __name__ == "__main__":
    unittest.main()
